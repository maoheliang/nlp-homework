import os
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from datetime import datetime
import pytesseract
from openpyxl import load_workbook
from docx.oxml.ns import qn
from imports import *


# Optional: for .doc to .docx conversion
import platform
if platform.system() == "Windows":
    import win32com.client
elif platform.system() in ["Linux", "Darwin"]:
    import subprocess

class IncrementalDocumentRetriever:
    def __init__(self, model_name="./models/bge-small-zh-v1.5", space="cosine", ef_construction=200, M=16, device="cpu"):
        self.model = SentenceTransformer(model_name, device=device)
        self.space = space
        self.ef_construction = ef_construction
        self.M = M
        self.index = None
        self.dim = None
        self.documents = []  # list of texts
        self.doc_ids = []    # list of unique int IDs
        self.id_counter = 0  # next doc id

    def _preprocess(self, text):
        return "为这个句子生成表示以用于检索相关文章：" + text

    def _get_embedding(self, text):
        processed = self._preprocess(text)
        return self.model.encode([processed], normalize_embeddings=True)[0]

    def _init_index(self, dim):
        self.index = hnswlib.Index(space=self.space, dim=dim)
        self.index.init_index(max_elements=10000, ef_construction=self.ef_construction, M=self.M)
        self.index.set_ef(50)

    def add_document(self, text: str):
        embedding = self._get_embedding(text)
        if self.index is None:
            self.dim = embedding.shape[0]
            self._init_index(self.dim)

        doc_id = self.id_counter
        self.index.add_items([embedding], [doc_id])
        self.documents.append(text)
        self.doc_ids.append(doc_id)
        self.id_counter += 1

    def search(self, query, top_k=5, return_scores=False) -> Union[List[Tuple[str, float]], str]:
        if self.index is None or self.id_counter == 0:
            raise ValueError("No documents indexed yet.")
        query_vec = self._get_embedding(query)
        labels, distances = self.index.knn_query(query_vec, k=min(top_k, self.id_counter))
        results = []
        for idx, dist in zip(labels[0], distances[0]):
            text = self.documents[self.doc_ids.index(idx)]
            score = 1 - dist
            results.append((text, score))
        if return_scores:
            return results
        else:
            return "\n".join([text for text, _ in results])
    
    def search_by_threshold(self, query: str, threshold: float = 0.3, return_scores: bool = False) -> Union[List[Tuple[str, float]], str]:
        if self.index is None or self.id_counter == 0:
            raise ValueError("No documents indexed yet.")
        
        query_vec = self._get_embedding(query)
        # brute-force 所有文档余弦相似度（比用 index 快不了太多，因为 hnswlib 不支持 range query）
        similarities = []
        for idx, doc_id in enumerate(self.doc_ids):
            doc_vec = self.index.get_items([doc_id])[0]
            score = np.dot(query_vec, doc_vec)  # 因为是 normalize 后的向量
            if score >= threshold:
                similarities.append((self.documents[idx], score))
        
        # 按相似度从高到低排序
        similarities.sort(key=lambda x: x[1], reverse=True)
        if return_scores:
            return similarities
        else:
            return "\n".join([text for text, _ in similarities])

    def save(self, index_path: str, meta_path: str):
        if self.index:
            self.index.save_index(index_path)
        meta = {
            "documents": self.documents,
            "doc_ids": self.doc_ids,
            "id_counter": self.id_counter,
            "dim": self.dim
        }
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)

    def load(self, index_path: str, meta_path: str):
        if not os.path.exists(index_path) or not os.path.exists(meta_path):
            raise FileNotFoundError("Index or metadata file not found.")

        with open(meta_path, "r", encoding="utf-8") as f:
            meta = json.load(f)
            self.documents = meta["documents"]
            self.doc_ids = meta["doc_ids"]
            self.id_counter = meta["id_counter"]
            self.dim = meta["dim"]

        self.index = hnswlib.Index(space=self.space, dim=self.dim)
        self.index.load_index(index_path)
        self.index.set_ef(50)

class DocumentProcessor:
    """
    处理文档的类
    """
    def __init__(self, use_retriever=False, model_name="./models/bge-small-zh-v1.5", device="cpu"):
        os.makedirs("./data", exist_ok=True)
        self.data = pd.DataFrame(columns=[
            "filename", "filetype", "content", "content_type", 
            "order_index", "timestamp"
        ])
        if use_retriever:
            self.retriever = IncrementalDocumentRetriever(model_name=model_name, device=device)

    def process_file(self, filepath):
        """
        处理文件
        :param filepath: 文件路径
        """
        # 获取文件扩展名
        ext = os.path.splitext(filepath)[1].lower()

        # 如果文件扩展名为.doc，则将其转换为.docx格式
        if ext == '.doc':
            filepath = self._convert_doc_to_docx(filepath)
            ext = '.docx'

        # 根据文件扩展名选择处理方法
        method = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".pptx": self._process_pptx,
            ".xlsx": self._process_excel,
            ".csv": self._process_csv
        }.get(ext)

        # 如果文件扩展名不在支持范围内，则抛出异常
        if not method:
            raise ValueError(f"Unsupported file type: {ext}")

        # 调用处理方法，获取文件内容
        contents = method(filepath)
        # 将文件路径、扩展名和内容添加到历史记录中
        self._append_to_data(filepath, ext, contents)

    def _append_to_data(self, filepath, filetype, contents):
        """
         将文件路径、扩展名和内容添加到历史记录中
        :param filepath: 文件路径
        :param filetype: 文件扩展名
        :param contents: 文件内容
        """
        # 获取当前时间戳
        now = datetime.now().isoformat()
        # 创建一个空列表，用于存储记录
        records = []
        # 遍历contents中的每个元素
        for order_index, content_type, content in contents:
            # 如果content不为空，则将记录添加到records列表中
            if content.strip():
                records.append({
                    "filename": os.path.abspath(os.path.normpath(filepath)),
                    "filetype": filetype,
                    "content": content.strip(),
                    "content_type": content_type,
                    "order_index": order_index,
                    "timestamp": now
                })
        # 如果records列表不为空，则将records转换为DataFrame，并添加到self.data中
        if records:
            new_records = pd.DataFrame(records)
            self.data = pd.concat([self.data, new_records], ignore_index=True)
            if hasattr(self, 'retriever'):
                record_str=f"====文件名：{filepath}====\n\n"
                for record in records:
                    record_str+=f"信息类型：{record['content_type']}\n\n"
                    record_str+=f"内容：{record['content']}\n"
                # print(record_str)
                self.retriever.add_document(record_str)

    def _process_pdf(self, filepath):
        """
         处理pdf文件
        """
        def clean_line_breaks(text: str) -> str:
            # 将非句末标点后的换行替换为空格
            pattern = r'(?<![。！？；：)”》）\]）])\n(?!\n)'  # 换行前不是结尾符，且后面不是段落空行
            return re.sub(pattern, '', text)
        contents = []
        full_text = []
        order = 1  # 表格从 1 开始，文本占用 order 0

        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    full_text.append(text)
                else:
                    try:
                        image = page.to_image(resolution=300).original
                        ocr_text = pytesseract.image_to_string(image)
                        full_text.append(ocr_text)
                    except Exception:
                        pass

                for table in page.extract_tables():
                    df = pd.DataFrame(table[1:], columns=table[0]) if len(table) > 1 else pd.DataFrame(table)
                    df = df.fillna("空值")
                    contents.append((order, 'table', df.to_csv(index=False)))
                    order += 1

        # 合并所有文本页为一条内容
        if full_text:
            merged_text = "\n".join(full_text)
            merged_text = clean_line_breaks(merged_text)
            contents.insert(0, (0, 'text', merged_text))

        return contents


    def _is_numbered_paragraph(self, paragraph):
        """
         判断段落是否为编号段
        """
        pPr = paragraph._p.find(qn('w:pPr'))
        if pPr is not None:
            numPr = pPr.find(qn('w:numPr'))
            return numPr is not None
        return False

    def _process_docx(self, filepath):
        """
         处理docx文件
        """
        doc = Document(filepath)
        contents = []
        order = 0

        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if self._is_numbered_paragraph(para):
                    text = "[编号段] " + text
                paragraphs.append(text)

        full_text = "\n".join(paragraphs)
        if full_text:
            contents.append((order, 'text', full_text))
            order += 1

        for table in doc.tables:
            data = [[cell.text if cell.text.strip() else "空值" for cell in row.cells] for row in table.rows]
            if data:
                header = data[0]
                rows = data[1:]
                df = pd.DataFrame(rows, columns=header)
            else:
                df = pd.DataFrame()
            contents.append((order, 'table', df.to_csv(index=False)))
            order += 1
        return contents

    def _process_pptx(self, filepath):
        """
         处理pptx文件
        """
        # 打开指定路径的pptx文件
        prs = Presentation(filepath)
        # 创建一个空列表，用于存储pptx文件中的内容
        contents = []
        # 初始化一个变量，用于记录内容的顺序
        order = 0
        # 遍历pptx文件中的每一张幻灯片
        for slide in prs.slides:
            # 遍历幻灯片中的每一个形状
            for shape in slide.shapes:
                # 如果形状有文本属性，并且文本不为空
                if hasattr(shape, "text") and shape.text.strip():
                    # 将内容的顺序、类型和文本添加到列表中
                    contents.append((order, 'text', shape.text))
                    # 顺序加1
                    order += 1
        # 返回存储内容的列表
        return contents
    
    def _process_excel(self, filepath):
        """
         处理excel文件
        """
        from openpyxl import load_workbook
        wb = load_workbook(filepath, data_only=False)
        contents = []
        order = 0
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=False):
                row_values = []
                for cell in row:
                    val = cell.value
                    if val is not None:
                        if cell.number_format and isinstance(val, (float, int)):
                            if '%' in cell.number_format:
                                val = f"{val:.2%}"
                            else:
                                val = str(val)
                        else:
                            val = str(val)
                    else:
                        val = "空值"
                    row_values.append(val)
                rows.append(row_values)
            if rows:
                header = rows[0]
                data = rows[1:]
                df = pd.DataFrame(data, columns=header)
            else:
                df = pd.DataFrame()
            df = df.fillna("空值")
            # csv_data = df.to_csv(index=False).replace('_x000D_', '\n')
            csv_data = df.to_csv(index=False)
            sheet_block = f"Sheet: {sheet_name}\n{csv_data}"
            contents.append((order, 'table', sheet_block))
            order += 1
        return contents

    # 定义一个函数，用于处理CSV文件
    def _process_csv(self, filepath):
        """
         处理CSV文件
        """
        # 使用pandas库读取CSV文件
        df = pd.read_csv(filepath)
        # 返回一个元组，包含0、'table'和df.to_csv(index=False)的值
        return [(0, 'table', df.to_csv(index=False))]

    def _convert_doc_to_docx(self, doc_path):
        """
         将doc文件转换为docx文件
        """
        # 判断操作系统类型
        if platform.system() == "Windows":
            # 使用win32com库调用Word应用程序
            word = win32com.client.Dispatch("Word.Application")
            # 打开指定路径的doc文件
            doc = word.Documents.Open(doc_path)
            # 定义新的文件路径，文件格式为docx
            new_path = doc_path + "x"  # .docx
            # 将doc文件另存为docx文件
            doc.SaveAs(new_path, FileFormat=16)  # wdFormatDocumentDefault
            # 关闭doc文件
            doc.Close()
            # 退出Word应用程序
            word.Quit()
            # 返回新的文件路径
            return new_path
        elif platform.system() in ["Linux", "Darwin"]:
            # 定义新的文件路径，文件格式为docx
            new_path = doc_path + "x"
            # 使用subprocess库调用libreoffice命令行工具将doc文件转换为docx文件
            try:
                subprocess.run(["libreoffice", "--headless", "--convert-to", "docx", "--outdir", os.path.dirname(doc_path), doc_path])
                # 返回新的文件路径
                return new_path
            except subprocess.CalledProcessError as e:
                print(f"Conversion failed: {e}")
        else:
            # 抛出异常，不支持当前操作系统
            raise EnvironmentError("Unsupported OS for .doc conversion")

    # 获取历史记录
    def get_data(self)-> pd.DataFrame:
        # 返回历史记录的副本
        return deepcopy(self.data)
    
    def get_str_data(self)-> str:
        """
        将数据转化为字符串表示输出
        """
        grouped = self.data.groupby("filename")
        str_data = ""
        for filename, group in grouped:
            str_data += f"\n=== 📄 文件: {filename}===\n\n"
            for _, row in group.iterrows():
                str_data += f"信息类型: {row['filetype']}\n\n"   
                str_data += f"内容: {row['content']}\n"
        return str_data


    def export_data(self, out_file="./data/report.xlsx"):
        """
         导出历史记录到Excel文件
         out_file: 导出的Excel文件路径
        """
        # 导入xlsxwriter库
        import xlsxwriter

        df = self.data.copy()

        # 清理内容中的异常回车字符，统一为 \n
        df['content'] = df['content'].astype(str)
        df['content'] = df['content'].str.replace('\r\n', '\n', regex=False)
        df['content'] = df['content'].str.replace('\r', '\n', regex=False)
        df['content'] = df['content'].str.replace('_x000D_', '\n', regex=False)

        with pd.ExcelWriter(out_file, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False, sheet_name='Report', startrow=1, header=False)

            workbook  = writer.book
            worksheet = writer.sheets['Report']

            # 写入表头
            for col_num, value in enumerate(df.columns):
                worksheet.write(0, col_num, value)

            # 自动换行格式
            wrap_format = workbook.add_format({'text_wrap': True, 'valign': 'top'})

            # 为每个单元格设置换行格式
            for row_num in range(1, len(df) + 1):
                for col_num in range(len(df.columns)):
                    worksheet.write(row_num, col_num, df.iat[row_num - 1, col_num], wrap_format)

            # 设置列宽（可调）
            worksheet.set_column(0, len(df.columns) - 1, width=40)

    def print_data_summary(self, max_lines=-1):
        """
        打印历史记录的摘要
        max_lines: 每个内容预览的最大行数，-1表示不限制
        """
        # 根据文件名对历史记录进行分组
        grouped = self.data.groupby("filename")
        # 遍历分组后的历史记录
        for filename, group in grouped:
            # 打印文件名和文件类型
            print(f"\n=== 📄 文件: {filename}（{group.iloc[0]['filetype']}）===")
            # 遍历每个文件的历史记录
            for _, row in group.iterrows():
                # 如果max_lines大于0，则截取内容的前max_lines行
                if max_lines > 0:
                    content_preview = row['content'].strip().splitlines()[:max_lines]
                # 否则，截取全部内容
                else:
                    content_preview = row['content'].strip().splitlines()
                # 将内容预览转换为字符串
                preview_str = "\n    ".join(content_preview)
                # 打印内容类型、顺序和时间
                print(f"  ▶ 类型: {row['content_type']}, 顺序: {row['order_index']}, 时间: {row['timestamp']}")
                # 打印内容预览
                print(f"    内容预览:\n    {preview_str}")



if __name__ == "__main__":
    processor = DocumentProcessor(use_retriever=True,device="cpu")
    processor.process_file("./data/2025年5月护理部理论知识培训.docx")
    processor.process_file("./data/2025年5月手卫生执行专项培训与评估总结.pdf")
    processor.process_file("./data/手卫生培训各科室参与与考核情况统计.xlsx")
    # processor.export_data()

    # processor.print_data_summary()